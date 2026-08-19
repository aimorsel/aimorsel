#!/usr/bin/env python3
"""多格式输入适配器（引擎策略第三层，P0-1）。

把 docx / xlsx / pptx / HTML 解析成与 opendataloader-pdf 完全相同的 JSON 结构树 schema，
下游管线（RAG 分块 / 表格导出 / QA / 合并 / MCP）不需要知道输入是什么格式；
图片则包装成单页 PDF 交给已有的 OCR 通道。

schema 约定（与底层 Java 引擎输出对齐，下游只依赖这些字段）：
- 根:       {"file name", "number of pages", "kids": [...]}
- heading:  {"type": "heading", "heading level": int, "page number": int, "content": str}
- paragraph:{"type": "paragraph", "page number": int, "content": str}
- list:     {"type": "list", "page number": int, "list items": [{"type": "list item", ...}]}
- table:    {"type": "table", "page number": int, "number of rows", "number of columns",
             "rows": [{"type": "table row", "row number", "cells": [{"type": "table cell",
             "row number", "column number", "row span", "column span", "content"}]}]}

「页」的语义：docx / HTML 无分页概念，全文算第 1 页；xlsx 每个工作表一页；pptx 每张幻灯片一页。

依赖全部可选（python-docx / openpyxl / python-pptx / Pillow），缺哪个只影响对应格式，
错误信息里直接给出 pip install 命令；HTML 用标准库 html.parser，零依赖。
"""

from __future__ import annotations

import re
from collections import Counter
from pathlib import Path

from i18n import tr

# 交给本模块解析的办公文档格式
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
# 网页（标准库 html.parser 解析，零依赖）
HTML_EXTENSIONS = {".html", ".htm", ".xhtml"}
# 全部走「解析成结构树」路径的格式（路由层用这个，别分别列举）
ADAPTER_EXTENSIONS = OFFICE_EXTENSIONS | HTML_EXTENSIONS
# 包装成 PDF 走 OCR 通道的图片格式
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp", ".gif"}


class AdapterError(RuntimeError):
    """适配器解析失败（含缺依赖），message 是给用户看的中文。"""


def _missing(package: str) -> AdapterError:
    return AdapterError(tr("缺少依赖 {package}，请运行: pip install {package}", package=package))


# ---------------------------------------------------------------- 结构树节点构造

def _document(path: Path, pages: int, kids: list[dict]) -> dict:
    return {"file name": path.name, "number of pages": max(pages, 1), "kids": kids}


def _heading(text: str, level: int, page: int) -> dict:
    return {"type": "heading", "heading level": max(1, min(int(level), 6)),
            "page number": page, "content": text}


def _paragraph(text: str, page: int) -> dict:
    return {"type": "paragraph", "page number": page, "content": text}


def _list(items: list[str], page: int) -> dict:
    return {"type": "list", "page number": page,
            "number of list items": len(items),
            "list items": [{"type": "list item", "page number": page, "content": t}
                           for t in items]}


def _table(grid: list[list[str]], page: int) -> dict:
    n_cols = max((len(r) for r in grid), default=0)
    rows = []
    for r, row in enumerate(grid, start=1):
        cells = [{"type": "table cell", "page number": page,
                  "row number": r, "column number": c, "row span": 1, "column span": 1,
                  "content": text}
                 for c, text in enumerate(row, start=1)]
        rows.append({"type": "table row", "row number": r, "cells": cells})
    return {"type": "table", "page number": page,
            "number of rows": len(grid), "number of columns": n_cols, "rows": rows}


# ---------------------------------------------------------------- 无标题文档的编号行提升（bench issue #12）

# 法条/规章类源文件常常一个 h 标签、一个标题样式都没有（日本 e-Gov 的 docx/html、EUR-Lex 的
# 指令正文），产物里一个层级都没有，下游分块无从下手。这里做一道**保守**启发式：
# 只有整篇零标题、且「短行 + 编号」命中 ≥3 处时才动，且只改命中的那几行——
# 有真实标题层级的文档一律不碰，宁可少给也不要造出假层级。
_PROMOTE_MIN_HITS = 3          # 少于这个数说明不是成体系的编号，可能只是正文里提了两句
_PROMOTE_MAX_CHARS = 60        # 「短行」上限：标题行不会长过这个数，条文正文一定超
# 命中后 **编号之后的部分**若含句末标点，就是「第五条 规定……。」这类正文引用，不提升
_SENTENCE_END = re.compile(r"[。．！？；]|[.!?;](?:\s|$)")
_CJK_NUM = r"[〇零一二三四五六七八九十百千万\d]+"
# (正则, 标题层级)：层级按「編/章 → 節 → 款/条」的常规法条层次给，与补出来的文档标题(1)错开
_NUMBERED_HEADING = (
    (re.compile(rf"^第\s*{_CJK_NUM}\s*[編篇部]"), 2),
    (re.compile(rf"^第\s*{_CJK_NUM}\s*章"), 2),
    (re.compile(rf"^第\s*{_CJK_NUM}\s*[節节]"), 3),
    (re.compile(rf"^第\s*{_CJK_NUM}\s*[款目]"), 4),
    (re.compile(rf"^第\s*{_CJK_NUM}\s*[条條]"), 4),
    # 多语言：EUR-Lex 同一份指令有 en/de/es/fr 四版，只认英文会让另外三版一个层级都拿不到。
    # 法语第一条写成序数词 "Article premier"，不单列的话每份 fr 正好少这一个标题
    (re.compile(r"^(?:part|chapter|title|chapitre|kapitel|cap[ií]tulo|t[ií]tulo|titel)\s+[\divxlcm]+\b", re.I), 2),
    (re.compile(r"^(?:section|sec\.|secci[óo]n|abschnitt)\s*\d", re.I), 3),
    (re.compile(r"^(?:article|art\.|artikel|art[ií]culo)\s*(?:\d|premi[eè]re?\b)", re.I), 4),
    (re.compile(r"^§+\s*\d"), 4),
)


def _match_numbered(text: str) -> tuple[int, bool] | None:
    """这行是不是「第三条」「Article 5」「§ 12」这类编号标题？是则给出 (层级, 编号后面是不是空的)。"""
    flat = " ".join(text.split())  # 全角空格 U+3000 也折成普通空格，只用于匹配，不改产物文本
    if not flat or len(flat) > _PROMOTE_MAX_CHARS:
        return None
    for pattern, level in _NUMBERED_HEADING:
        m = pattern.match(flat)
        if m:
            rest = flat[m.end():]
            return None if _SENTENCE_END.search(rest) else (level, not rest.strip(" .,，、-–—:："))
    return None


def _numbered_heading_level(text: str) -> int | None:
    matched = _match_numbered(text)
    return matched[0] if matched else None


def _promote_numbered_headings(doc: dict) -> int:
    """整篇没有标题层级时，把顶层的编号短行段落提升成标题，返回提升的行数。"""
    kids = doc.get("kids") or []
    heads = [i for i, k in enumerate(kids) if k.get("type") == "heading"]
    # 允许「只有一个开头的一级标题」——那是 <title>/文件名补出来的文档标题，不算层级
    if heads and (len(heads) > 1 or heads[0] != 0 or kids[0].get("heading level") != 1):
        return 0
    hits = [(i, m[0], m[1], " ".join((k.get("content") or "").split())) for i, k in enumerate(kids)
            if k.get("type") == "paragraph" and (m := _match_numbered(k.get("content") or ""))]
    # 反复出现同一行的是页眉/页脚（SEC 财报每页都印一次 "PART I"，一份里 100 多次），不是标题。
    # 带标题文字的（「第二章　最低賃金」）允许出现两次——法条的「目次」会先列一遍，正文里再出现一次；
    # 光秃秃只有编号的（"PART III"）出现两次就当页眉，正文标题不会这么写
    seen = Counter(t for _, _, _, t in hits)
    hits = [(i, level) for i, level, bare, t in hits if seen[t] < (2 if bare else 3)]
    if len(hits) < _PROMOTE_MIN_HITS:
        return 0
    for i, level in hits:
        kids[i] = _heading(kids[i]["content"], level, kids[i].get("page number", 1))
    return len(hits)


# ---------------------------------------------------------------- 各格式解析

def parse_office(path: Path) -> dict:
    """把 docx/xlsx/pptx/HTML 解析成 JSON 结构树。失败抛 AdapterError（中文信息）。

    名字沿用「office」是历史原因；HTML 加进来后语义是「所有适配器格式」（`parse_document` 同义）。
    """
    suffix = path.suffix.lower()
    parsers = {".docx": _parse_docx, ".xlsx": _parse_xlsx, ".pptx": _parse_pptx}
    parser = parsers.get(suffix) or (_parse_html if suffix in HTML_EXTENSIONS else None)
    try:
        if parser is not None:
            doc = parser(path)
            _promote_numbered_headings(doc)  # 零标题文档的兜底层级（bench issue #12）
            return doc
    except AdapterError:
        raise
    except Exception as err:  # 解析库抛的各种异常，统一翻译成可读错误
        raise AdapterError(tr("解析 {name} 失败：{err}", name=path.name, err=err)) from err
    raise AdapterError(tr("不支持的格式：{suffix}", suffix=suffix))


parse_document = parse_office


def _parse_docx(path: Path) -> dict:
    try:
        import docx
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError:
        raise _missing("python-docx") from None

    doc = docx.Document(str(path))
    kids: list[dict] = []
    pending_list: list[str] = []

    def flush_list() -> None:
        if pending_list:
            kids.append(_list(list(pending_list), page=1))
            pending_list.clear()

    # doc.paragraphs 和 doc.tables 是分开的两个列表，要保持阅读顺序必须走 body 的 XML 子节点
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "") if para.style is not None else ""
            level = _docx_heading_level(style)
            if level:
                flush_list()
                kids.append(_heading(text, level, page=1))
            elif "list" in style.lower():
                pending_list.append(text)
            else:
                flush_list()
                kids.append(_paragraph(text, page=1))
        elif isinstance(child, CT_Tbl):
            flush_list()
            table = Table(child, doc)
            grid = [["\n".join(p.text for p in cell.paragraphs).strip()
                     for cell in row.cells]
                    for row in table.rows]
            if grid:
                kids.append(_table(grid, page=1))
    flush_list()
    return _document(path, 1, kids)


def _docx_heading_level(style_name: str) -> int | None:
    """Word 样式名 -> 标题级别。"Heading 2"/"heading 2" -> 2，"Title" -> 1，其余 None。"""
    lowered = style_name.lower()
    if lowered == "title":
        return 1
    if lowered.startswith("heading"):
        tail = lowered[len("heading"):].strip()
        try:
            return int(tail)
        except ValueError:
            return 1
    return None


def _parse_xlsx(path: Path) -> dict:
    try:
        import openpyxl
    except ImportError:
        raise _missing("openpyxl") from None

    # data_only=True 读公式的缓存计算值而不是公式本身
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        kids: list[dict] = []
        pages = 0
        for page, sheet in enumerate(wb.worksheets, start=1):
            pages = page
            kids.append(_heading(sheet.title, 1, page))
            grid = [[_cell_str(v) for v in row] for row in sheet.iter_rows(values_only=True)]
            while grid and all(c == "" for c in grid[-1]):  # 去掉尾部全空行
                grid.pop()
            n_cols = max((len(r) for r in grid), default=0)
            while n_cols and all(len(r) < n_cols or r[n_cols - 1] == "" for r in grid):
                n_cols -= 1  # 去掉尾部全空列
            grid = [r[:n_cols] + [""] * (n_cols - len(r[:n_cols])) for r in grid]
            if grid and n_cols:
                kids.append(_table(grid, page))
        return _document(path, pages, kids)
    finally:
        wb.close()


def _cell_str(value) -> str:
    """单元格值 -> 文本。整数值的浮点去掉 .0，日期用 ISO 格式。

    openpyxl 把日期单元格一律读成 datetime（零点），显示时裁掉多余的 00:00:00。
    """
    import datetime

    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    if isinstance(value, datetime.datetime):
        if value.time() == datetime.time.min:
            return value.date().isoformat()
        return value.isoformat(sep=" ")
    if hasattr(value, "isoformat"):  # date / time
        return value.isoformat()
    return str(value)


def _parse_pptx(path: Path) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        raise _missing("python-pptx") from None

    prs = Presentation(str(path))
    kids: list[dict] = []
    pages = 0
    for page, slide in enumerate(prs.slides, start=1):
        pages = page
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                grid = [[cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows]
                if grid:
                    kids.append(_table(grid, page))
            elif getattr(shape, "has_text_frame", False):
                is_title = title_shape is not None and shape == title_shape
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip() or para.text.strip()
                    if not text:
                        continue
                    if is_title:
                        kids.append(_heading(text, 2, page))  # 幻灯片标题统一二级，文件名占一级
                        is_title = False  # 标题框里的后续段落按正文处理
                    else:
                        kids.append(_paragraph(text, page))
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                kids.append(_paragraph(tr("（演讲者备注）") + notes, page))
    return _document(path, pages, kids)



# ---------------------------------------------------------------- HTML -> 结构树

# 这些元素的内容整段丢弃（脚本/样式/隐藏模板/矢量图；<head> 里只留 <title>）
# <math> 不在此列：MathML 公式另行捕获成 LaTeX 文本（见 _HtmlTreeBuilder 的 math 分支，bench issue #1）
# ix:header / ix:hidden 是 XBRL 内联标注的元数据容器（SEC 财报），浏览器不渲染，
# 进正文就是一千多字符的 "0000732712 2025 FY FALSE us-gaap:CommonStockMember…"（bench issue #11）
_HTML_SKIP = {"script", "style", "noscript", "template", "svg", "iframe", "object",
              "ix:header", "ix:hidden"}
# 行内样式声明的不可见元素：我们不做 CSS 级联，只认写在 style 属性里的这两条
# （SEC 财报把 XBRL 事实塞在 <div style="display:none"> 里，一份就有三千多处）
_CSS_HIDDEN = re.compile(r"(?:^|;)\s*(?:display\s*:\s*none|visibility\s*:\s*hidden)\s*(?:!\s*important\s*)?(?:;|$)", re.I)


# 空元素（void）没有结束标签：`<img style="display:none">` 若按普通元素起 skip，
# 就永远等不到 </img>，后文全被吞掉（与 bench/truth/from_html.py 的 VOID 处理保持一致）
_HTML_VOID = {"area", "base", "basefont", "br", "col", "embed", "frame", "hr", "img", "input",
              "keygen", "link", "meta", "param", "source", "track", "wbr"}


def _html_hidden(attrs: dict) -> bool:
    """这个元素在浏览器里根本不显示 → 内容不进正文。"""
    return "hidden" in attrs or bool(_CSS_HIDDEN.search(attrs.get("style") or ""))


def _clean_tex(tex: str) -> str:
    """去掉 MediaWiki 一类生成器包在 LaTeX 外面的 {\\displaystyle ...} / {\\textstyle ...} 壳。

    只有首个 "{" 的配对 "}" 恰好是末字符时才剥（"{\\displaystyle a}{\\displaystyle b}" 这类并列壳原样保留，
    否则会剥出花括号不配对的残缺 LaTeX）。
    """
    tex = " ".join(tex.split())
    m = re.match(r"\{\s*\\(?:displaystyle|textstyle|scriptstyle)\b\s*", tex)
    if not m or not tex.endswith("}"):
        return tex
    depth = 0
    for i, ch in enumerate(tex):
        if ch == "\\":
            continue
        if ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                return tex[m.end():-1].strip() if i == len(tex) - 1 else tex
    return tex


def _math_markup(tex: str, display: bool) -> str:
    """公式的文本形式：行内 $...$，独立公式 $$...$$ 自成一行。"""
    return f"\n$${tex}$$\n" if display else f"${tex}$"
# 块级元素：开始/结束时把攒着的行内文本收成一个段落（或 li / 单元格内容）
_HTML_BLOCK = {
    "p", "div", "section", "article", "header", "footer", "main", "aside", "nav",
    "blockquote", "pre", "figure", "figcaption", "hr", "address", "details", "summary",
    "form", "fieldset", "dl", "dt", "dd", "center", "body", "html",
}
_HTML_HEADINGS = {"h1": 1, "h2": 2, "h3": 3, "h4": 4, "h5": 5, "h6": 6}
# HTML 允许省略这些元素的结束标签：`<li>甲<li>乙</ul>` 里第一个 li 没有 </li>，html.parser
# **不会**替我们补。被丢弃的元素若正好是它们之一（`<td style="display:none">`），只等同名
# 结束标签就永远等不到，后文整篇被吞（实测 kids 为空、无报错）。键 = 可省略结束标签的元素，
# 值 = 出现哪些开始标签就意味着它已经闭合
_HTML_IMPLIED_END = {
    "p": _HTML_BLOCK | set(_HTML_HEADINGS) | {"ul", "ol", "menu", "li", "table", "tr", "td", "th"},
    "li": {"li"},
    "td": {"td", "th", "tr"}, "th": {"td", "th", "tr"}, "tr": {"tr"},
    "option": {"option", "optgroup"}, "optgroup": {"optgroup"},
    "dt": {"dt", "dd"}, "dd": {"dt", "dd"},
    "thead": {"tbody", "tfoot"}, "tbody": {"tbody", "tfoot"}, "tfoot": {"tbody"},
}


def _sniff_html_encoding(raw: bytes) -> str | None:
    """从 BOM / <meta charset> / <?xml encoding> 里猜声明的编码。"""
    if raw.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    head = raw[:4096]
    m = (re.search(rb"""<meta[^>]+charset\s*=\s*["']?\s*([-\w.:]+)""", head, re.I)
         or re.search(rb"""<\?xml[^>]+encoding\s*=\s*["']([-\w.:]+)""", head, re.I))
    return m.group(1).decode("ascii", "ignore") if m else None


def _read_html_text(path: Path) -> str:
    """按声明编码 → utf-8 → utf-8 宽松 三级读取。"""
    raw = path.read_bytes()
    declared = _sniff_html_encoding(raw)
    for enc in (declared, "utf-8"):
        if not enc:
            continue
        try:
            return raw.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="replace")


class _HtmlTreeBuilder:
    """事件驱动的 HTML → 结构树转换（标准库 html.parser，容忍未闭合标签）。

    规则：h1-h6 → heading；ul/ol 的 li → list（嵌套列表拍平进外层，保持阅读顺序）；
    table 的 tr/td/th → table（<th> 与 <td> 同等对待，caption 成独立段落）；
    其余块级元素边界处把行内文本收成 paragraph；<br> 保留为换行；<pre> 内保留空白。
    """

    def __init__(self) -> None:
        from html.parser import HTMLParser

        builder = self

        class _P(HTMLParser):
            def handle_starttag(self, tag, attrs):
                builder._start(tag, dict(attrs))

            def handle_endtag(self, tag):
                builder._end(tag)

            def handle_data(self, data):
                builder._data(data)

        self._parser = _P(convert_charrefs=True)
        self.kids: list[dict] = []
        self.title = ""
        self._buf: list[str] = []
        self._skip = 0             # 是否正处在被丢弃的元素里
        self._skip_tag = ""        # 触发丢弃的标签名
        self._skip_stack: list[str] = []   # 被丢弃区域内部还开着的元素（用于识别它到底在哪结束）
        self._pre = 0
        self._in_title = False
        self._heading: int | None = None
        self._lists: list[dict] = []    # {"items": [...], "open": bool}
        self._tables: list[dict] = []   # {"rows": [...], "row": list|None, "cell": list|None}
        # <math> 捕获：优先 <annotation encoding="…tex"> 的 LaTeX，其次 alttext 属性，最后 MathML 词元拼接
        self._math = 0                  # <math> 同名嵌套深度
        self._math_display = False
        self._math_alt = ""
        self._math_tex: list[str] = []
        self._math_tokens: list[str] = []
        self._math_ann = 0              # 0=不在 annotation 里；1=TeX annotation；2=其他 annotation(-xml)，内容丢弃
        self._last_math = ""            # 刚写出的公式（去重 MediaWiki 紧随其后的 alt=LaTeX 的回退图片）

    # ---- 对外
    def feed(self, text: str) -> None:
        self._parser.feed(text)
        self._parser.close()
        self._flush()
        while self._tables:
            self._close_table()
        while self._lists:
            self._close_list()

    # ---- 文本缓冲
    def _data(self, data: str) -> None:
        if self._skip:
            return
        if self._math:
            if self._math_ann == 1:
                self._math_tex.append(data)
            elif self._math_ann == 0:
                self._math_tokens.append(data)
            return
        if self._in_title:
            self.title += data
            return
        if not self._pre:
            data = data.replace("\r", " ").replace("\n", " ").replace("\t", " ")
        if data:
            if data.strip():
                self._last_math = ""  # 去重只针对紧跟在公式后面的回退图片
            self._buf.append(data)

    def _take(self) -> str:
        raw = "".join(self._buf)
        self._buf.clear()
        if self._pre:  # 保留缩进，只去掉首尾空行
            return "\n".join(ln.rstrip() for ln in raw.split("\n")).strip("\n")
        lines = [re.sub(r"[ \t\f\v\xa0]+", " ", ln).strip() for ln in raw.split("\n")]
        return "\n".join(ln for ln in lines if ln)

    def _flush(self, heading: int | None = None) -> None:
        """把缓冲文本按当前上下文归位：单元格 > 列表项 > 标题 > 段落。"""
        text = self._take()
        if not text:
            return
        table = self._tables[-1] if self._tables else None
        if table is not None and table["cell"] is not None:
            table["cell"].append(text)
            return
        lst = self._lists[-1] if self._lists else None
        if lst is not None and lst["open"]:
            lst["items"].append(text)
            lst["open"] = False  # 一个 li 内的多段文本合并成一个项：后续文本追加到该项
            lst["append_to_last"] = True
            return
        if lst is not None and lst.get("append_to_last") and lst["items"]:
            lst["items"][-1] += " " + text
            return
        if heading:
            self.kids.append(_heading(text, heading, page=1))
        else:
            self.kids.append(_paragraph(text, page=1))

    def _emit_math(self) -> None:
        """</math> 收口：把公式以文本形式写进行内缓冲。"""
        tex = _clean_tex("".join(self._math_tex)) or _clean_tex(self._math_alt) or "".join(self._math_tokens)
        tex = " ".join(tex.split())
        self._math_tex, self._math_tokens, self._math_alt = [], [], ""
        if not tex:
            return
        in_paragraph = not self._heading and not (self._lists and self._lists[-1]["open"]) \
            and not (self._tables and self._tables[-1]["cell"] is not None)
        self._buf.append(_math_markup(tex, self._math_display and in_paragraph))  # 标题/列表项/单元格里只能行内
        self._last_math = tex

    # ---- 标签事件
    def _start(self, tag: str, attrs: dict) -> None:
        if self._skip:
            # 内部那些省略了结束标签的元素，遇到后继标签就算闭合
            while self._skip_stack and tag in _HTML_IMPLIED_END.get(self._skip_stack[-1], ()):
                self._skip_stack.pop()
            if not self._skip_stack and tag in _HTML_IMPLIED_END.get(self._skip_tag, ()):
                self._end_skip()          # 被丢弃的元素自己省略了结束标签：就地收口，别吞后文
                self._start(tag, attrs)
                return
            if tag not in _HTML_VOID:
                self._skip_stack.append(tag)
            return
        if self._math:
            if tag == "math":
                self._math += 1
            elif tag == "annotation":
                self._math_ann = 1 if "tex" in (attrs.get("encoding") or "").lower() else 2
            elif tag == "annotation-xml":
                self._math_ann = 2
            elif tag in _HTML_BLOCK or tag in _HTML_HEADINGS or tag in ("li", "table", "tr", "td", "th", "ul", "ol"):
                # 未闭合的 <math> 遇到块级标签：就地收口，别把后文全吞进公式
                self._math = 0
                self._emit_math()
                self._start(tag, attrs)
            return
        if tag in _HTML_SKIP or _html_hidden(attrs) or (tag != "img" and (attrs.get("aria-hidden") or "").lower() == "true"):
            # aria-hidden 的非图片元素是给屏幕阅读器藏起来的视觉副本（KaTeX 的 katex-html、MathJax v2 的
            # <nobr>）：MathML 已写出公式，视觉副本再进正文就重复了。回退图片（MediaWiki）也带 aria-hidden，
            # 但它的 alt 是唯一来源之一，交给下面的 img 分支去重处理
            if tag in _HTML_VOID:
                return  # 无结束标签，起了 skip 就再也关不掉，直接丢弃这个标签本身即可
            self._skip, self._skip_tag, self._skip_stack = 1, tag, []
            return
        if tag not in ("math", "img", "br"):
            self._last_math = ""
        if tag == "math":
            self._math = 1
            self._math_display = (attrs.get("display") or "").lower() == "block"
            self._math_alt = attrs.get("alttext") or ""
            self._math_tex, self._math_tokens, self._math_ann = [], [], 0
            return
        if tag == "img":
            # 只有回退图片的公式（MediaWiki 无 MathML 输出、KaTeX 图片模式等）：alt 里是 LaTeX
            alt = attrs.get("alt") or ""
            if alt and "math" in (attrs.get("class") or "").lower():
                tex = _clean_tex(alt)
                if tex and tex != self._last_math:
                    self._buf.append(_math_markup(tex, False))
                    self._last_math = tex
            return
        if tag == "title":
            self._in_title = True
            return
        if tag == "br":
            self._buf.append("\n")
            return
        if tag in _HTML_HEADINGS:
            self._flush()
            self._heading = _HTML_HEADINGS[tag]
            return
        if tag in ("ul", "ol", "menu"):
            self._flush()
            self._lists.append({"items": [], "open": False, "append_to_last": False})
            return
        if tag == "li":
            self._flush()
            if not self._lists:  # 游离的 li，当作单项列表
                self._lists.append({"items": [], "open": False, "append_to_last": False})
            self._lists[-1]["open"] = True
            self._lists[-1]["append_to_last"] = False
            return
        if tag == "table":
            self._flush()
            self._tables.append({"rows": [], "row": None, "cell": None})
            return
        if tag == "tr" and self._tables:
            self._flush()
            self._end_row()
            self._tables[-1]["row"] = []
            return
        if tag in ("td", "th") and self._tables:
            self._flush()
            table = self._tables[-1]
            if table["row"] is None:
                table["row"] = []
            self._end_cell()
            table["cell"] = []
            return
        if tag == "caption" and self._tables:
            self._flush()
            return
        if tag == "pre":
            self._flush()
            self._pre += 1
            return
        if tag in _HTML_BLOCK:
            self._flush()

    def _end_skip(self) -> None:
        self._skip, self._skip_tag, self._skip_stack = 0, "", []

    def _end(self, tag: str) -> None:
        if self._skip:
            if tag in self._skip_stack:                 # 内部元素闭合（中间那些是省略了结束标签的）
                while self._skip_stack.pop() != tag:
                    pass
            elif tag == self._skip_tag:                 # 自己闭合；内部没闭合的元素一并作废
                self._end_skip()
            elif self._skip_tag in _HTML_IMPLIED_END:
                # 祖先的结束标签先到（`<ul><li style=display:none>x</ul>`）：被丢弃的元素
                # 省略了自己的结束标签，就地收口再把这个结束标签交回去
                self._end_skip()
                self._end(tag)
            return
        if self._math:
            if tag in ("annotation", "annotation-xml"):
                self._math_ann = 0
            elif tag == "math":
                self._math -= 1
                if not self._math:
                    self._emit_math()
            return
        if tag == "title":
            self._in_title = False
            return
        if tag in _HTML_HEADINGS:
            level, self._heading = self._heading, None
            self._flush(heading=level or _HTML_HEADINGS[tag])
            return
        if tag in ("ul", "ol", "menu"):
            self._flush()
            self._close_list()
            return
        if tag == "li":
            self._flush()
            if self._lists:
                self._lists[-1]["open"] = False
                self._lists[-1]["append_to_last"] = False
            return
        if tag == "table":
            self._flush()
            self._close_table()
            return
        if tag == "tr" and self._tables:
            self._flush()
            self._end_row()
            return
        if tag in ("td", "th") and self._tables:
            self._flush()
            self._end_cell()
            return
        if tag == "caption" and self._tables:
            # caption 文本落在段落里，放在表格之前
            text = self._take()
            if text:
                self.kids.append(_paragraph(text, page=1))
            return
        if tag == "pre":
            self._flush()
            self._pre = max(0, self._pre - 1)
            return
        if tag in _HTML_BLOCK:
            self._flush()

    # ---- 列表 / 表格收尾
    def _close_list(self) -> None:
        lst = self._lists.pop()
        items = [i for i in lst["items"] if i]
        if not items:
            return
        if self._lists:  # 嵌套列表：拍平进外层，保持顺序
            self._lists[-1]["items"].extend(items)
            self._lists[-1]["append_to_last"] = False
        else:
            self.kids.append(_list(items, page=1))

    def _end_cell(self) -> None:
        table = self._tables[-1]
        if table["cell"] is not None:
            table["row"].append(" ".join(table["cell"]).strip())
            table["cell"] = None

    def _end_row(self) -> None:
        table = self._tables[-1]
        self._end_cell()
        if table["row"] is not None:
            if any(c for c in table["row"]):
                table["rows"].append(table["row"])
            table["row"] = None

    def _close_table(self) -> None:
        self._end_row()
        table = self._tables.pop()
        rows = table["rows"]
        if not rows:
            return
        n_cols = max(len(r) for r in rows)
        grid = [r + [""] * (n_cols - len(r)) for r in rows]
        if self._tables and self._tables[-1]["cell"] is not None:
            # 表中表：把内表拍成文本塞进外层单元格（罕见，不追求结构还原）
            self._tables[-1]["cell"].append(" ".join(c for r in grid for c in r if c))
        else:
            self.kids.append(_table(grid, page=1))


def _parse_html(path: Path) -> dict:
    """HTML → 结构树（全文算第 1 页）。<title> 只在正文没有一级标题时补作首个标题。"""
    builder = _HtmlTreeBuilder()
    builder.feed(_read_html_text(path))
    kids = builder.kids
    title = " ".join(builder.title.split())
    if title and not any(k["type"] == "heading" and k["heading level"] == 1 for k in kids):
        kids.insert(0, _heading(title, 1, page=1))
    return _document(path, 1, kids)


# ---------------------------------------------------------------- 图片 -> 单页 PDF

# ---- 倾斜校正（bench issue #8）
# 真实扫描件几乎都带一两度倾斜，版面分析按倾斜的行切块会把内容切碎（实测 ±3° 的 tiff 只保留
# 46–49% 文本，同页 docling 是我们的 2.2 倍）。投影法：把图缩小成灰度图，在 ±5° 里试着转正，
# 取「每行亮度的方差最大」的角度——文字行摆平时行与行间白，方差最大。纯 Pillow，无 numpy。
DESKEW_MAX_DEG = 5.0       # 只纠正轻微倾斜；再大就不是扫描歪了，是页面本身横放，别乱转
# 阈值 1.0 是实测定的：0.5–0.8 度的真实扫描件（xfund 五份）转正后 char_sim **全部下降**
# （-0.01 ~ -0.10，重采样的模糊比那点倾斜更伤 OCR），1.4 度以上则一致大涨。宁可不转
DESKEW_MIN_DEG = 1.0
_DESKEW_PROBE_PX = 800     # 估角用的缩略图长边；再大只是变慢，角度分辨率并不提高


def detect_skew(img) -> float:
    """估计图片的倾斜角（度）。返回值是**要转多少度才摆正**，逆时针为正。"""
    from PIL import Image

    gray = img.convert("L")
    scale = _DESKEW_PROBE_PX / max(gray.size)
    if scale < 1:
        gray = gray.resize((max(1, int(gray.width * scale)), max(1, int(gray.height * scale))),
                           Image.BILINEAR)

    def score(angle: float) -> float:
        rot = gray.rotate(angle, resample=Image.BILINEAR, fillcolor=255) if angle else gray
        # 旋转在四角留下白色楔形，会盖过文字信号：只取中间 80% 算投影
        w, h = rot.size
        rot = rot.crop((int(w * 0.1), int(h * 0.1), int(w * 0.9), int(h * 0.9)))
        if rot.width < 1 or rot.height < 2:
            return 0.0   # 极细长的图（1×1 占位图、横幅条）裁完中间 80% 什么都不剩，没有行投影可算
        rows = list(rot.resize((1, rot.height), Image.BOX).getdata())  # 每行的平均亮度
        if len(rows) < 2:
            return 0.0
        mean = sum(rows) / len(rows)
        return sum((v - mean) ** 2 for v in rows) / len(rows)

    def best(candidates: list[float]) -> float:
        return max(candidates, key=score)

    step = 10  # 0.1 度为单位算，避免浮点累加误差
    coarse = best([a / step for a in range(int(-DESKEW_MAX_DEG * step), int(DESKEW_MAX_DEG * step) + 1, 5)])
    fine = best([coarse + d / step for d in range(-4, 5)])
    return fine if abs(fine) <= DESKEW_MAX_DEG else 0.0


def deskew_image(img):
    """轻微倾斜的图转正；角度不足 DESKEW_MIN_DEG 时原样返回（返回 (图, 实际转的角度)）。"""
    from PIL import Image

    angle = detect_skew(img)
    if abs(angle) < DESKEW_MIN_DEG:
        return img, 0.0
    return img.rotate(angle, resample=Image.BICUBIC, expand=True, fillcolor=(255, 255, 255)), angle


def image_to_pdf(image_path: Path, pdf_path: Path, deskew: bool = True) -> float:
    """把图片包装成 PDF（多帧 TIFF/GIF 每帧一页），供已有的 OCR 通道处理。

    `deskew` 打开时顺带把轻微倾斜的页面转正，返回转动角度绝对值的最大值（没转就是 0）。
    """
    try:
        from PIL import Image
    except ImportError:
        raise _missing("pillow") from None

    try:
        turned = 0.0
        with Image.open(image_path) as img:
            frames = []
            for i in range(getattr(img, "n_frames", 1)):
                img.seek(i)
                frame = img.convert("RGB")
                if deskew:
                    try:
                        frame, angle = deskew_image(frame)
                    except Exception:   # 校正只是增强：失败就用原图，别把能转的文件变成失败
                        angle = 0.0
                    turned = max(turned, abs(angle))
                frames.append(frame)
            first, rest = frames[0], frames[1:]
            first.save(pdf_path, "PDF", save_all=bool(rest), append_images=rest)
        return turned
    except AdapterError:
        raise
    except Exception as err:
        raise AdapterError(tr("读取图片 {name} 失败：{err}", name=image_path.name, err=err)) from err


# ---------------------------------------------------------------- 结构树 -> 各格式渲染

def _table_markdown(node: dict) -> str:
    lines: list[str] = []
    for i, row in enumerate(node.get("rows") or []):
        cells = [(cell.get("content") or "").replace("|", "\\|").replace("\n", " ")
                 for cell in row.get("cells") or []]
        lines.append("| " + " | ".join(cells) + " |")
        if i == 0:
            lines.append("|" + " --- |" * len(cells))
    return "\n".join(lines)


def _render_blocks(tree: dict, page_separator: str | None, heading_prefix: bool) -> str:
    """markdown / 纯文本共用的线性渲染。page_separator 含 %page-number% 占位符。"""
    parts: list[str] = []
    last_page = None
    for node in tree.get("kids") or []:
        page = node.get("page number")
        if page_separator and last_page is not None and page != last_page:
            parts.append(page_separator.replace("%page-number%", str(page)).strip("\n"))
        last_page = page
        node_type = node.get("type")
        if node_type == "heading":
            text = node.get("content") or ""
            level = node.get("heading level") or 1
            parts.append(("#" * level + " " + text) if heading_prefix else text)
        elif node_type == "table":
            parts.append(_table_markdown(node))
        elif node_type == "list":
            parts.append("\n".join("- " + (item.get("content") or "")
                                   for item in node.get("list items") or []))
        else:
            text = node.get("content") or ""
            if text:
                parts.append(text)
    return "\n\n".join(p for p in parts if p.strip()) + "\n"


def render_markdown(tree: dict, page_separator: str | None = None) -> str:
    return _render_blocks(tree, page_separator, heading_prefix=True)


def render_text(tree: dict, page_separator: str | None = None) -> str:
    return _render_blocks(tree, page_separator, heading_prefix=False)


def render_html(tree: dict) -> str:
    import html as html_mod

    body: list[str] = []
    for node in tree.get("kids") or []:
        node_type = node.get("type")
        if node_type == "heading":
            level = min(node.get("heading level") or 1, 6)
            body.append(f"<h{level}>{html_mod.escape(node.get('content') or '')}</h{level}>")
        elif node_type == "table":
            rows_html = []
            for row in node.get("rows") or []:
                cells = "".join(f"<td>{html_mod.escape(cell.get('content') or '')}</td>"
                                for cell in row.get("cells") or [])
                rows_html.append(f"<tr>{cells}</tr>")
            body.append("<table>" + "".join(rows_html) + "</table>")
        elif node_type == "list":
            items = "".join(f"<li>{html_mod.escape(item.get('content') or '')}</li>"
                            for item in node.get("list items") or [])
            body.append(f"<ul>{items}</ul>")
        else:
            text = node.get("content") or ""
            if text:
                body.append(f"<p>{html_mod.escape(text)}</p>")
    title = html_mod.escape(tree.get("file name") or "")
    return ("<!DOCTYPE html>\n<html>\n<head>\n<meta charset=\"utf-8\">\n"
            f"<title>{title}</title>\n</head>\n<body>\n" + "\n".join(body) + "\n</body>\n</html>\n")
