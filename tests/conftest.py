"""pytest 公共 fixture：样例文件生成 + Java 可用性探测。

测试 PDF 是提交入库的固定 fixture（tests/data/simple.pdf，reportlab 生成，
两页多级标题）——不依赖 cupsfilter（仅 macOS）也不依赖网络。
office/图片样例用运行时依赖现场生成，与被测代码用同一批库，无额外依赖。
"""

from __future__ import annotations

import os
import shutil
import sys
from pathlib import Path

import pytest

# 测试断言全部基于中文文案，锁死语言，避免 CI 机器 locale（en）改变输出
os.environ.setdefault("MORSEL_LANG", "zh")

PROJECT = Path(__file__).resolve().parent.parent
if str(PROJECT) not in sys.path:
    sys.path.insert(0, str(PROJECT))

DATA = Path(__file__).parent / "data"

JAVA_AVAILABLE = shutil.which("java") is not None
requires_java = pytest.mark.skipif(not JAVA_AVAILABLE, reason="系统没有 java，跳过真实引擎集成测试")


@pytest.fixture(scope="session")
def text_pdf() -> Path:
    """两页英文文本 PDF（标题 Annual Report / Introduction / Financial Details）。"""
    return DATA / "simple.pdf"


@pytest.fixture(scope="session")
def scanned_pdf(tmp_path_factory) -> Path:
    """纯图片 PDF（无文字层），模拟扫描件。"""
    from PIL import Image

    path = tmp_path_factory.mktemp("scan") / "scanned.pdf"
    Image.new("RGB", (600, 800), (200, 30, 30)).save(path, "PDF")
    return path


@pytest.fixture(scope="session")
def sample_docx(tmp_path_factory) -> Path:
    import docx

    doc = docx.Document()
    doc.add_heading("年度报告", level=1)
    doc.add_paragraph("这是第一段正文，介绍整体情况。")
    doc.add_heading("财务部分", level=2)
    doc.add_paragraph("苹果", style="List Bullet")
    doc.add_paragraph("香蕉", style="List Bullet")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "项目"
    table.rows[0].cells[1].text = "金额"
    table.rows[1].cells[0].text = "收入"
    table.rows[1].cells[1].text = "100"
    doc.add_paragraph("表格后的结尾段落。")
    path = tmp_path_factory.mktemp("docx") / "report.docx"
    doc.save(path)
    return path


@pytest.fixture(scope="session")
def sample_xlsx(tmp_path_factory) -> Path:
    import datetime

    import openpyxl

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "销售"
    ws1.append(["月份", "销量", "日期"])
    ws1.append(["一月", 42, datetime.date(2026, 1, 31)])
    ws1.append(["二月", 3.5, None])
    ws2 = wb.create_sheet("成本")
    ws2.append(["项目", "成本"])
    ws2.append(["原料", 7])
    path = tmp_path_factory.mktemp("xlsx") / "book.xlsx"
    wb.save(path)
    return path


@pytest.fixture(scope="session")
def sample_pptx(tmp_path_factory) -> Path:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "项目启动"
    slide.placeholders[1].text = "目标一\n目标二"
    slide.notes_slide.notes_text_frame.text = "记得强调预算"
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "里程碑"
    path = tmp_path_factory.mktemp("pptx") / "deck.pptx"
    prs.save(path)
    return path


@pytest.fixture(scope="session")
def sample_png(tmp_path_factory) -> Path:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (600, 200), "white")
    ImageDraw.Draw(img).text((30, 80), "Hello AImorsel", fill="black")
    path = tmp_path_factory.mktemp("png") / "photo.png"
    img.save(path)
    return path
