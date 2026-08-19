#!/usr/bin/env python3
"""自造语料（bench/PLAN.md §1.2「自造图片集/Office 集」+ 冒烟集），全部带 A 档精确真值。

    python -m bench.make_synthetic [--out bench/corpus/synthetic] [--seed 0]

产出：PDF（reportlab，多语种、带标题/表格/列表）、docx/xlsx/pptx、PNG/JPG/TIFF 图片
（PIL 渲染，供 OCR）、HTML；每份旁边一个 ``<id>.truth.json``：
``{"text", "headings", "tables", "paragraphs"}``，并追加到 ``manifest.jsonl``。
字体：macOS 自带 Arial Unicode（覆盖 CJK/西里尔/希腊，**不做阿拉伯连字**——阿语靠真实语料 UN ODS）。
"""
from __future__ import annotations

import argparse
import json
import random
from pathlib import Path

BENCH = Path(__file__).resolve().parent
FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/Library/Fonts/Arial Unicode.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
]

# 语料库：lang -> {domain: (标题, [段落...])}。句子是真实语义的短文，不是 lorem。
BANK: dict[str, dict[str, tuple[str, list[str]]]] = {
    "zh": {
        "it": ("分布式系统的一致性模型", [
            "在分布式系统中，一致性模型定义了并发读写操作可见性的规则。线性一致性要求每个操作看起来都在某个瞬间原子地生效。",
            "最终一致性放松了这一要求：只要停止写入，所有副本最终会收敛到同一状态。这种模型换来了更高的可用性和更低的延迟。",
            "工程实践中常用的折中是因果一致性，它保证有因果关系的操作被所有节点以相同顺序观察到。",
        ]),
        "math": ("线性代数中的特征值分解", [
            "对于 n 阶方阵 A，若存在非零向量 v 和标量 λ 使得 Av = λv，则称 λ 为 A 的特征值，v 为对应的特征向量。",
            "当 A 有 n 个线性无关的特征向量时，A 可以对角化：A = PDP⁻¹，其中 D 是由特征值构成的对角矩阵。",
            "实对称矩阵总可以正交对角化，这是谱定理的内容，也是主成分分析的理论基础。",
        ]),
        "law": ("合同的成立与生效", [
            "依照民法典规定，当事人订立合同可以采用书面形式、口头形式或者其他形式。承诺生效时合同成立。",
            "依法成立的合同自成立时生效，但法律另有规定或者当事人另有约定的除外。附条件的合同自条件成就时生效。",
            "无民事行为能力人实施的民事法律行为无效；限制民事行为能力人订立的纯获利益的合同有效。",
        ]),
        "business": ("季度经营情况说明", [
            "本季度公司实现营业收入十二点五亿元，同比增长百分之十八，其中海外业务贡献了主要增量。",
            "毛利率维持在百分之四十一，较上季度提升一点二个百分点，主要得益于原材料价格回落与产品结构优化。",
            "经营性现金流净额为二点三亿元，资产负债率下降至百分之三十五，财务结构进一步稳健。",
        ]),
        "edu": ("小学阶段阅读能力培养", [
            "阅读能力的培养应从兴趣入手，低年级以绘本和短篇故事为主，逐步过渡到章节书。",
            "教师应引导学生在阅读中提问、预测和复述，这三种策略被证明能显著提升理解水平。",
            "家庭共读每周不少于三次，每次二十分钟，长期坚持对词汇量增长的作用最为明显。",
        ]),
    },
    "en": {
        "it": ("Consistency Models in Distributed Systems", [
            "A consistency model specifies the rules under which concurrent reads and writes become visible. Linearizability requires every operation to appear to take effect atomically at some instant.",
            "Eventual consistency relaxes this: once writes stop, all replicas converge to the same state. The trade-off buys higher availability and lower latency.",
            "Causal consistency is the common middle ground, guaranteeing that causally related operations are observed in the same order by every node.",
        ]),
        "math": ("Eigenvalue Decomposition", [
            "For a square matrix A of order n, if a nonzero vector v and scalar λ satisfy Av = λv, then λ is an eigenvalue and v an eigenvector.",
            "When A has n linearly independent eigenvectors it is diagonalizable: A = PDP⁻¹ where D is the diagonal matrix of eigenvalues.",
            "Real symmetric matrices are always orthogonally diagonalizable; this spectral theorem underlies principal component analysis.",
        ]),
        "law": ("Formation of Contracts", [
            "A contract is formed when an offer is accepted, supported by consideration, and the parties intend to create legal relations.",
            "Acceptance must be communicated to the offeror; silence does not generally constitute acceptance under common law.",
            "Contracts entered into by minors are voidable at the minor's option, except for contracts for necessaries.",
        ]),
        "business": ("Quarterly Operating Review", [
            "Revenue for the quarter reached 1.25 billion, up 18 percent year on year, with international business contributing most of the growth.",
            "Gross margin held at 41 percent, up 1.2 points from the prior quarter, driven by lower input costs and a richer product mix.",
            "Operating cash flow was 230 million and the debt-to-asset ratio fell to 35 percent.",
        ]),
        "edu": ("Building Reading Skills in Primary School", [
            "Reading instruction should begin with interest: picture books and short stories in early grades, moving gradually to chapter books.",
            "Teachers should prompt students to question, predict and retell while reading; these three strategies measurably improve comprehension.",
            "Shared family reading at least three times a week for twenty minutes has the clearest long-term effect on vocabulary growth.",
        ]),
    },
    "es": {
        "it": ("Modelos de consistencia en sistemas distribuidos", [
            "Un modelo de consistencia define las reglas bajo las cuales las lecturas y escrituras concurrentes se hacen visibles.",
            "La consistencia eventual relaja este requisito: cuando cesan las escrituras, todas las réplicas convergen al mismo estado.",
            "La consistencia causal garantiza que las operaciones relacionadas causalmente se observen en el mismo orden en todos los nodos.",
        ]),
        "law": ("Formación del contrato", [
            "El contrato se perfecciona por el mero consentimiento, y desde entonces obliga al cumplimiento de lo expresamente pactado.",
            "El consentimiento se manifiesta por el concurso de la oferta y de la aceptación sobre la cosa y la causa.",
            "Será nulo el consentimiento prestado por error, violencia, intimidación o dolo.",
        ]),
        "edu": ("Comprensión lectora en primaria", [
            "La enseñanza de la lectura debe partir del interés del alumno, con álbumes ilustrados en los primeros cursos.",
            "Preguntar, predecir y resumir durante la lectura mejora de forma medible la comprensión.",
            "La lectura compartida en familia tres veces por semana tiene el efecto más claro sobre el vocabulario.",
        ]),
    },
    "de": {
        "math": ("Eigenwertzerlegung", [
            "Für eine quadratische Matrix A heißt λ Eigenwert, wenn ein Vektor v ≠ 0 mit Av = λv existiert.",
            "Besitzt A n linear unabhängige Eigenvektoren, so ist A diagonalisierbar: A = PDP⁻¹.",
            "Reelle symmetrische Matrizen sind stets orthogonal diagonalisierbar (Spektralsatz).",
        ]),
        "law": ("Zustandekommen eines Vertrags", [
            "Ein Vertrag kommt durch zwei übereinstimmende Willenserklärungen zustande: Angebot und Annahme.",
            "Die Annahme muss dem Antragenden zugehen; Schweigen gilt grundsätzlich nicht als Annahme.",
            "Verträge beschränkt Geschäftsfähiger bedürfen der Einwilligung des gesetzlichen Vertreters.",
        ]),
        "business": ("Quartalsbericht", [
            "Der Umsatz stieg im Quartal um 18 Prozent auf 1,25 Milliarden, getragen vom Auslandsgeschäft.",
            "Die Bruttomarge lag bei 41 Prozent, 1,2 Punkte über dem Vorquartal.",
            "Der operative Cashflow betrug 230 Millionen; die Verschuldungsquote sank auf 35 Prozent.",
        ]),
    },
    "fr": {
        "it": ("Modèles de cohérence", [
            "Un modèle de cohérence précise quand les écritures concurrentes deviennent visibles aux lecteurs.",
            "La cohérence à terme garantit la convergence des réplicas une fois les écritures arrêtées.",
            "La cohérence causale préserve l'ordre des opérations liées par une relation de cause à effet.",
        ]),
        "edu": ("Apprentissage de la lecture", [
            "L'apprentissage de la lecture doit partir de l'intérêt de l'élève, avec des albums illustrés au début.",
            "Questionner, prédire et reformuler pendant la lecture améliore nettement la compréhension.",
            "La lecture partagée en famille trois fois par semaine a l'effet le plus net sur le vocabulaire.",
        ]),
    },
    "ja": {
        "it": ("分散システムにおける一貫性モデル", [
            "一貫性モデルは、並行する読み書き操作がいつ他のノードから見えるようになるかの規則を定める。",
            "結果整合性では、書き込みが止まればすべてのレプリカが最終的に同じ状態に収束する。",
            "因果一貫性は、因果関係のある操作がすべてのノードで同じ順序で観測されることを保証する。",
        ]),
        "business": ("四半期業績の概況", [
            "当四半期の売上高は前年同期比18パーセント増の12億5千万となり、海外事業が成長を牽引した。",
            "売上総利益率は41パーセントと前四半期から1.2ポイント改善した。",
            "営業キャッシュフローは2億3千万、負債比率は35パーセントに低下した。",
        ]),
    },
    "ru": {
        "math": ("Спектральное разложение", [
            "Число λ называется собственным значением матрицы A, если существует ненулевой вектор v такой, что Av = λv.",
            "Если у A есть n линейно независимых собственных векторов, то A диагонализуема: A = PDP⁻¹.",
            "Вещественные симметричные матрицы всегда ортогонально диагонализуемы.",
        ]),
    },
}

TABLE_BANK = {
    "zh": [["指标", "本季度", "上季度"], ["营业收入", "12.5 亿", "10.6 亿"], ["毛利率", "41%", "39.8%"], ["现金流", "2.3 亿", "1.9 亿"]],
    "default": [["Metric", "Q2", "Q1"], ["Revenue", "1.25 B", "1.06 B"], ["Gross margin", "41%", "39.8%"], ["Cash flow", "230 M", "190 M"]],
}
LIST_BANK = {"zh": ["第一项要点", "第二项要点", "第三项要点"], "default": ["First point", "Second point", "Third point"]}


def font_path() -> str:
    for p in FONT_CANDIDATES:
        if Path(p).exists():
            return p
    raise SystemExit("找不到覆盖 CJK 的 TTF 字体，请在 FONT_CANDIDATES 里加一个")


def build_doc(lang: str, domain: str, rng: random.Random, with_table=True, with_list=True) -> dict:
    """生成文档逻辑结构：sections=[{"heading","paras"}]，可选表格/列表。返回结构 + 真值。"""
    title, paras = BANK[lang][domain]
    sub_titles = {
        "zh": ["背景", "方法", "结论"], "ja": ["背景", "方法", "結論"],
    }.get(lang, ["Background", "Method", "Conclusion"])
    sections = [{"heading": sub_titles[i], "paras": [p]} for i, p in enumerate(paras)]
    table = (TABLE_BANK.get(lang) or TABLE_BANK["default"]) if with_table else None
    items = (LIST_BANK.get(lang) or LIST_BANK["default"]) if with_list else None
    text_parts, headings, paragraphs = [title], [title], []
    for s in sections:
        headings.append(s["heading"])
        text_parts.append(s["heading"])
        for p in s["paras"]:
            text_parts.append(p)
            paragraphs.append(p)
    if table:
        text_parts += [" ".join(r) for r in table]
    if items:
        text_parts += items
    truth = {"text": " ".join(text_parts), "headings": headings, "paragraphs": paragraphs,
             "tables": [table] if table else []}
    return {"title": title, "sections": sections, "table": table, "items": items, "truth": truth}


# ---------- 各格式写出 ----------

def write_pdf(doc: dict, path: Path, font: str, pages: int = 1) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
    from reportlab.lib import colors

    if "U" not in pdfmetrics.getRegisteredFontNames():
        pdfmetrics.registerFont(TTFont("U", font))
    h1 = ParagraphStyle("h1", fontName="U", fontSize=18, leading=24, spaceAfter=10)
    h2 = ParagraphStyle("h2", fontName="U", fontSize=14, leading=18, spaceBefore=8, spaceAfter=6)
    body = ParagraphStyle("b", fontName="U", fontSize=10.5, leading=15, spaceAfter=6)
    flow = [Paragraph(doc["title"], h1)]
    for i, s in enumerate(doc["sections"]):
        flow.append(Paragraph(s["heading"], h2))
        flow += [Paragraph(p, body) for p in s["paras"]]
        if pages > 1 and i == 0:
            flow.append(PageBreak())
    if doc["table"]:
        t = Table(doc["table"])
        t.setStyle(TableStyle([("FONTNAME", (0, 0), (-1, -1), "U"), ("GRID", (0, 0), (-1, -1), 0.5, colors.black)]))
        flow += [Spacer(1, 8), t]
    if doc["items"]:
        flow += [Paragraph("• " + it, body) for it in doc["items"]]
    SimpleDocTemplate(str(path), pagesize=A4).build(flow)


def write_docx(doc: dict, path: Path) -> None:
    import docx

    d = docx.Document()
    d.add_heading(doc["title"], level=1)
    for s in doc["sections"]:
        d.add_heading(s["heading"], level=2)
        for p in s["paras"]:
            d.add_paragraph(p)
    if doc["table"]:
        t = d.add_table(rows=len(doc["table"]), cols=len(doc["table"][0]))
        for i, r in enumerate(doc["table"]):
            for j, c in enumerate(r):
                t.cell(i, j).text = c
    if doc["items"]:
        for it in doc["items"]:
            d.add_paragraph(it, style="List Bullet")
    d.save(str(path))


def write_xlsx(doc: dict, path: Path) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append([doc["title"]])
    for r in doc["table"] or []:
        ws.append(r)
    ws2 = wb.create_sheet("Notes")
    for s in doc["sections"]:
        ws2.append([s["heading"]])
        for p in s["paras"]:
            ws2.append([p])
    wb.save(str(path))


def write_pptx(doc: dict, path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = doc["title"]
    for s in doc["sections"]:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = s["heading"]
        sl.placeholders[1].text = "\n".join(s["paras"])
    if doc["table"]:
        sl = prs.slides.add_slide(prs.slide_layouts[5])
        rows, cols = len(doc["table"]), len(doc["table"][0])
        from pptx.util import Inches
        shape = sl.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3))
        for i, r in enumerate(doc["table"]):
            for j, c in enumerate(r):
                shape.table.cell(i, j).text = c
    prs.save(str(path))


def write_html(doc: dict, path: Path) -> None:
    import html as H

    parts = [f"<h1>{H.escape(doc['title'])}</h1>"]
    for s in doc["sections"]:
        parts.append(f"<h2>{H.escape(s['heading'])}</h2>")
        parts += [f"<p>{H.escape(p)}</p>" for p in s["paras"]]
    if doc["table"]:
        rows = "".join("<tr>" + "".join(f"<td>{H.escape(c)}</td>" for c in r) + "</tr>" for r in doc["table"])
        parts.append(f"<table>{rows}</table>")
    if doc["items"]:
        parts.append("<ul>" + "".join(f"<li>{H.escape(i)}</li>" for i in doc["items"]) + "</ul>")
    path.write_text("<!doctype html><html><head><meta charset='utf-8'></head><body>" + "\n".join(parts) + "</body></html>", encoding="utf-8")


def write_image(doc: dict, path: Path, font: str, size_px: int = 30, noise: bool = False, rotate: float = 0.0) -> None:
    """PIL 渲染成图片（供 OCR）。字号 ≥ 28px 才在 OCR 质量边界内（实测更小的字会整行丢失）。"""
    from PIL import Image, ImageDraw, ImageFilter, ImageFont

    W = 1654  # A4 @ 200dpi
    f_h1 = ImageFont.truetype(font, int(size_px * 1.6))
    f_h2 = ImageFont.truetype(font, int(size_px * 1.25))
    f_b = ImageFont.truetype(font, size_px)
    lines: list[tuple[str, object]] = [(doc["title"], f_h1)]
    for s in doc["sections"]:
        lines.append((s["heading"], f_h2))
        for p in s["paras"]:
            # 简单折行：按像素宽
            buf = ""
            for ch in p:
                if f_b.getlength(buf + ch) > W - 200:
                    lines.append((buf, f_b))
                    buf = ""
                buf += ch
            lines.append((buf, f_b))
    for r in doc["table"] or []:
        lines.append(("    ".join(r), f_b))
    for it in doc["items"] or []:
        lines.append(("• " + it, f_b))
    H_total = 100 + sum(int(f.size * 1.5) for _, f in lines) + 100
    img = Image.new("RGB", (W, max(H_total, 800)), "white")
    d = ImageDraw.Draw(img)
    y = 100
    for text, f in lines:
        d.text((100, y), text, font=f, fill="black")
        y += int(f.size * 1.5)
    if rotate:
        img = img.rotate(rotate, expand=True, fillcolor="white")
    if noise:
        img = img.filter(ImageFilter.GaussianBlur(0.8))
    suf = path.suffix.lower()
    if suf in (".jpg", ".jpeg"):
        img.save(str(path), quality=70)
    elif suf in (".tif", ".tiff"):
        img.convert("1").save(str(path), compression="group4")  # 二值传真压缩，模拟真实扫描件
    else:
        img.save(str(path))


# ---------- 生成计划 ----------

def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=str(BENCH / "corpus" / "synthetic"))
    ap.add_argument("--seed", type=int, default=0)
    a = ap.parse_args()
    out = Path(a.out)
    out.mkdir(parents=True, exist_ok=True)
    font = font_path()
    rng = random.Random(a.seed)
    manifest = out.parent / "manifest.jsonl"
    rows = []

    def add(doc_id, path, fmt, lang, domain, truth, layout, note=""):
        (out / f"{doc_id}.truth.json").write_text(json.dumps(truth, ensure_ascii=False, indent=1), encoding="utf-8")
        rows.append({"id": doc_id, "path": str(path.relative_to(out.parent)), "format": fmt, "lang": lang,
                     "domain": domain, "layout": layout, "source": "synthetic", "license": "CC0",
                     "truth_type": "A", "truth": str((out / f"{doc_id}.truth.json").relative_to(out.parent)),
                     "note": note})

    for lang, domains in BANK.items():
        for domain in domains:
            doc = build_doc(lang, domain, rng)
            base = f"syn_{lang}_{domain}"
            # PDF：单页 + 双页
            for pages in (1, 2):
                p = out / f"{base}_p{pages}.pdf"
                write_pdf(doc, p, font, pages=pages)
                add(f"{base}_pdf{pages}", p, "pdf", lang, domain, doc["truth"], "single-column")
            # Office
            p = out / f"{base}.docx"; write_docx(doc, p); add(f"{base}_docx", p, "docx", lang, domain, doc["truth"], "office")
            p = out / f"{base}.pptx"; write_pptx(doc, p); add(f"{base}_pptx", p, "pptx", lang, domain, doc["truth"], "office")
            xl_truth = {"text": " ".join([doc["title"]] + [" ".join(r) for r in doc["table"]] +
                                         [x for s in doc["sections"] for x in [s["heading"]] + s["paras"]]),
                        "tables": [doc["table"]]}
            p = out / f"{base}.xlsx"; write_xlsx(doc, p); add(f"{base}_xlsx", p, "xlsx", lang, domain, xl_truth, "office")
            # HTML
            p = out / f"{base}.html"; write_html(doc, p); add(f"{base}_html", p, "html", lang, domain, doc["truth"], "web")
            # 图片：干净 PNG / 压缩 JPG / 轻微倾斜模糊 TIFF
            p = out / f"{base}_clean.png"; write_image(doc, p, font); add(f"{base}_png", p, "png", lang, domain, doc["truth"], "scan-clean")
            p = out / f"{base}_q70.jpg"; write_image(doc, p, font, noise=True); add(f"{base}_jpg", p, "jpg", lang, domain, doc["truth"], "scan-noisy")
            p = out / f"{base}_skew.tiff"; write_image(doc, p, font, rotate=1.5); add(f"{base}_tiff", p, "tiff", lang, domain, doc["truth"], "scan-skew")

    # 追加/覆盖 manifest 里 source=synthetic 的行
    old = []
    if manifest.exists():
        old = [l for l in manifest.read_text(encoding="utf-8").splitlines() if l.strip() and '"source": "synthetic"' not in l]
    manifest.write_text("\n".join(old + [json.dumps(r, ensure_ascii=False) for r in rows]) + "\n", encoding="utf-8")
    print(f"生成 {len(rows)} 份合成语料 → {out}，manifest 共 {len(old) + len(rows)} 行")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
