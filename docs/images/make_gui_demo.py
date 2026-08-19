"""录桌面 GUI（`morsel gui`）的演示动图：勾选项 → 点 Convert → 进度条 / 日志推进 → 完成汇总。

抓帧照 make_gui_png.py：起 GUI、用 root.after 排定时器摆状态，`screencapture -x -l <窗口id>`
按窗口截图（需要 macOS 屏幕录制权限）。鼠标光标与点击波纹照 make_web_demo.py：
截图不含鼠标，光标是 PIL 事后画上去的——**不移动真实鼠标**，录制期间不抢用户的机器。
原生文件对话框截不到，所以文件用程序塞进列表，光标滑到 "Add files" 示意一下。

先造演示文件（列表会显示真实路径，所以放 /tmp 下，别让个人目录进 README；
英文版文件名与 gui-en.png 同一套）：
    python docs/images/make_gui_demo.py --make-demo-files

然后（英文版必须 MORSEL_LANG=en；帧落在 frames_gui/，GIF 直接写到目标路径）：
    rm -rf /tmp/aimorsel-demo/output
    MORSEL_LANG=en python docs/images/make_gui_demo.py docs/images/gui-demo-en.gif

转换是真的在跑（PDF / xlsx / pptx 走引擎，png 走 OCR 服务——服务在线时这一段最好看），
帧间隔按真实墙钟时间记录，重复帧由 PIL 自动合并。
"""
import os
import subprocess
import sys
import time
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))  # 仓库根
os.environ.setdefault("MORSEL_LANG", "zh")

from PIL import Image, ImageDraw

EN = os.environ["MORSEL_LANG"] == "en"
DEMO = Path("/tmp/aimorsel-demo/docs" if EN else "/tmp/aimorsel-demo/文档")
OUTPUT = Path("/tmp/aimorsel-demo/output")
FRAMES = Path("frames_gui")
SCALE = 1            # 截图是 Retina 2x，GIF 缩到逻辑像素（窗口本身已经 860 宽，别再缩）
COLORS = 128         # GIF 调色板大小（全帧共用）


def make_demo_files() -> None:
    """造四份小演示文件（PDF / 图片 / xlsx / pptx），内容是泛泛的示例数据。"""
    DEMO.mkdir(parents=True, exist_ok=True)
    names = (("quarterly-report.pdf", "annual-figures.xlsx", "contract-scan.png", "slides.pptx") if EN
             else ("季度报告.pdf", "年报.xlsx", "合同扫描件.png", "课件.pptx"))
    pdf, xlsx, png, pptx = (DEMO / n for n in names)

    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(str(pdf), pagesize=A4)
    c.setFont("Helvetica-Bold", 18); c.drawString(72, 780, "Quarterly Report")
    c.setFont("Helvetica", 12)
    for i, line in enumerate(["Revenue grew 12% quarter over quarter.",
                              "Operating margin: 18.4%", "Headcount: 142",
                              "", "Region    Q1      Q2", "North     1,204   1,388",
                              "South       932   1,051"]):
        c.drawString(72, 740 - 18 * i, line)
    c.save()

    import openpyxl
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Summary"
    ws.append(["Year", "Revenue", "Profit"])
    for y, r, p in ((2022, 8.1, 1.2), (2023, 9.4, 1.6), (2024, 10.9, 2.0)):
        ws.append([y, r, p])
    wb.save(xlsx)

    im = Image.new("RGB", (1240, 600), "white"); d = ImageDraw.Draw(im)
    from PIL import ImageFont
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 36)
    except Exception:
        font = ImageFont.load_default()
    for i, line in enumerate(["SERVICE AGREEMENT", "", "This agreement is made between",
                              "the Client and the Provider.", "Term: 12 months",
                              "Monthly fee: 1,250.00"]):
        d.text((80, 60 + 70 * i), line, fill="black", font=font)
    im.save(png)

    from pptx import Presentation
    prs = Presentation()
    for title, body in (("Product Overview", "Local-first document conversion"),
                        ("Roadmap", "Q3: OCR\nQ4: MCP server")):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title; s.placeholders[1].text = body
    prs.save(pptx)
    print("演示文件已写入", DEMO)


if "--make-demo-files" in sys.argv:
    make_demo_files()
    sys.exit(0)

OUT = Path(sys.argv[1] if len(sys.argv) > 1 else "gui-demo.gif")
FRAMES.mkdir(exist_ok=True)
for old in FRAMES.glob("*.png"):
    old.unlink()

import tkinter as tk
from tkinter import messagebox, ttk

from aimorsel import morsel_gui
from aimorsel.i18n import tr

# 完成弹窗会挡住 mainloop（而且截不到），录制时只要日志里的汇总行
messagebox.showinfo = lambda *a, **k: None
messagebox.showwarning = lambda *a, **k: None

try:
    from tkinterdnd2 import TkinterDnD
    root = TkinterDnD.Tk()
except Exception:
    root = tk.Tk()

app = morsel_gui.ConverterApp(root)
files = sorted(DEMO.iterdir()) if DEMO.is_dir() else []
if len(files) < 4:
    sys.exit(f"演示文件不够：{DEMO}（先跑 --make-demo-files）")
app.output_dir = OUTPUT
app.output_label.config(text=str(app.output_dir))
# 同 make_gui_png.py：收窄两个可伸缩区，让 RAG / 表格导出 / 质量自检那两行进得来
app.listbox.config(height=4)
app.log_text.config(height=6)
root.update_idletasks()
root.geometry(f"{860 if EN else 760}x{min(root.winfo_reqheight(), 900)}")
root.update_idletasks(); root.update()
root.lift(); root.attributes("-topmost", True); root.update()
# 从终端起的 Tk 窗口多半不是 key window，而且录制期间用户随手切个 App 它就失焦，macOS 会把
# 复选框/进度条画成灰色的「非活动」样式（gui-en.png 里是蓝的，动图里一灰就像被禁用）。
# 不去抢焦点（osascript 激活进程能让它变蓝，但用户一切回 Chrome 就又灰了，反复抢等于抢机器）；
# 改成每 100ms 把所有 ttk 控件的 `background` 状态位清掉——Tk 靠这个位画非活动样式，
# 清掉就按活动样式画，窗口本身焦点在哪无所谓。实测清掉后下一次重绘即恢复蓝色。
def keep_active_look():
    stack = [root]
    while stack:
        w = stack.pop()
        if isinstance(w, ttk.Widget):
            try:
                w.state(["!background"])
            except tk.TclError:
                pass
        stack.extend(w.winfo_children())
    root.after(100, keep_active_look)


keep_active_look()
root.update()

# 起录前确认语言生效：标题和 Convert 按钮都得是目标语言，否则白跑
assert app.convert_btn.cget("text") == tr("开始转换"), app.convert_btn.cget("text")
print("标题:", root.title(), "| 按钮:", app.convert_btn.cget("text"))


def find_widget(cls, text):
    stack = [root]
    while stack:
        w = stack.pop()
        if isinstance(w, cls) and str(w.cget("text")) == text:
            return w
        stack.extend(w.winfo_children())
    raise LookupError(text)


BTN_ADD = find_widget(ttk.Button, tr("添加文件"))
CHECKS = [(find_widget(ttk.Checkbutton, tr(t)), getattr(app, v)) for t, v in
          (("RAG 分块（输出 chunks.jsonl）", "rag_var"), ("表格导出 CSV", "tables_var"), ("质量自检", "qa_var"))]

# ------------------------------------------------------------------ 抓帧 + 光标
cursor = [640, 120]      # 当前光标位置（窗口客户区逻辑坐标）
CLICK_RING = [0]
frames = []              # [(path, 墙钟时刻)]


def widget_center(w, dx=None):
    """控件中心 → 截图坐标（截图 = 窗口客户区，不含标题栏）。"""
    x = w.winfo_rootx() - root.winfo_rootx() + (dx if dx is not None else w.winfo_width() / 2)
    y = w.winfo_rooty() - root.winfo_rooty() + w.winfo_height() / 2
    return [round(x), round(y)]


def draw_cursor(im):
    d = ImageDraw.Draw(im)
    x, y = cursor
    if CLICK_RING[0]:
        r = 10 + 6 * (3 - CLICK_RING[0])
        d.ellipse([x - r, y - r, x + r, y + r], outline="white", width=6)
        d.ellipse([x - r, y - r, x + r, y + r], outline=(31, 95, 208), width=3)
        CLICK_RING[0] -= 1
    arrow = [(x, y), (x, y + 17), (x + 4, y + 13), (x + 7, y + 20),
             (x + 10, y + 19), (x + 7, y + 12), (x + 12, y + 12)]
    d.polygon(arrow, fill="white", outline="black")


def shot(tag):
    p = FRAMES / f"{len(frames):03d}_{tag}.png"
    # `-l <winfo_id>` 在 tk 上不可用（winfo_id 是 NSView 指针，不是 CGWindowID，报
    # could not create image from window）；按窗口客户区的屏幕矩形截（-R），不含标题栏。
    # 窗口已置顶，矩形里就是它自己。
    region = f"{root.winfo_rootx()},{root.winfo_rooty()},{root.winfo_width()},{root.winfo_height()}"
    r = subprocess.run(["screencapture", "-x", "-o", "-R", region, str(p)], capture_output=True)
    if r.returncode != 0 or not p.exists():
        raise SystemExit("screencapture 失败（检查屏幕录制权限）")
    im = Image.open(p).convert("RGB")
    logical = (root.winfo_width() * SCALE, root.winfo_height() * SCALE)
    if im.size != logical:                      # Retina 截图是 2x，缩到逻辑像素
        im = im.resize(logical, Image.LANCZOS)
    draw_cursor(im)
    im.save(p)
    frames.append((p, time.monotonic()))


def glide(to, steps=5):
    x0, y0 = cursor
    for i in range(1, steps + 1):
        cursor[0] = round(x0 + (to[0] - x0) * i / steps)
        cursor[1] = round(y0 + (to[1] - y0) * i / steps)
        shot("move"); yield 70


def click(w=None, dx=None):
    CLICK_RING[0] = 3
    for _ in range(2):
        shot("click"); yield 120


# ------------------------------------------------------------------ 演示流程（生成器，yield 毫秒）
done = [False]
_orig_on_done = app._on_done


def _on_done(summary):
    _orig_on_done(summary)
    done[0] = True


app._on_done = _on_done


def script():
    cursor[:] = widget_center(app.listbox)          # 光标起手停在空列表里
    for _ in range(4):
        shot("idle"); yield 300

    yield from glide(widget_center(BTN_ADD))        # 示意「Add files」（对话框截不到）
    yield from click()
    app.pdfs.extend(files); app._refresh_list()
    app._log(tr("已添加 {n} 个文件", n=len(files)) if not EN else f"Added {len(files)} files (PDF / image / xlsx / pptx)")
    for _ in range(3):
        shot("picked"); yield 350

    for w, var in CHECKS:                            # 勾三个场景选项
        yield from glide(widget_center(w, dx=9))
        yield from click()
        var.set(True)
        shot("checked"); yield 250

    cx, cy = widget_center(app.convert_btn)          # 点在按钮右下一点，波纹别盖住 "Convert" 字样
    yield from glide([cx + 22, cy + 5])
    yield from click()
    app._start_conversion()                          # 真转换，后台线程
    t0 = time.monotonic()
    while not done[0] and time.monotonic() - t0 < 180:
        shot("run"); yield 250                       # 进度条 / 日志逐行推进——动图的价值所在，别快进
    for _ in range(4):                               # 停在完成汇总
        shot("done"); yield 300


steps = script()


def tick():
    try:
        delay = next(steps)
    except StopIteration:
        root.quit(); return
    root.after(delay, tick)


root.after(800, tick)
root.mainloop()
try:
    root.destroy()
except Exception:
    pass

# ------------------------------------------------------------------ 拼 GIF
# 体积要点（实测 55 帧 860×855）：① 全帧共用一张调色板（几帧拼图 quantize 出来）且不抖动，
# 否则每帧各自量化出的噪点让帧间差异全是满幅；② disposal=1（保留上一帧、只写变化区域），
# PIL 对 disposal=2 会把每帧整幅写出 → 5–7 MB，减帧/降色都救不回来；disposal=1 后不到 1 MB。
rgb = [Image.open(p).convert("RGB") for p, _ in frames]
sample = [rgb[0], rgb[len(rgb) // 3], rgb[2 * len(rgb) // 3], rgb[-1]]
montage = Image.new("RGB", (sample[0].width, sum(s.height for s in sample)))
for i, s in enumerate(sample):
    montage.paste(s, (0, i * s.height))
palette = montage.quantize(colors=COLORS)
imgs = [im.quantize(palette=palette, dither=Image.Dither.NONE) for im in rgb]
durs = []
for i, (p, t) in enumerate(frames):
    nxt = frames[i + 1][1] if i + 1 < len(frames) else t + 1.2   # 末帧多留一秒
    durs.append(int(max(60, min(1500, (nxt - t) * 1000))))
imgs[0].save(OUT, save_all=True, append_images=imgs[1:], duration=durs, loop=0, optimize=True, disposal=1)
g = Image.open(OUT)
print(f"帧 {len(frames)} 张 → GIF {g.n_frames} 帧，{g.size}，{OUT.stat().st_size / 1024:.0f} KB")
